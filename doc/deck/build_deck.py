"""헬름 IRP세븐서티액티브 발표자료(PPTX) 생성기.

투자제안서 v1.3 과 **같은 한 벌의 숫자**(_deck/metrics.json 재현치)를 쓰고,
과제 브리프(KDT AI퀀트 07)의 평가 배점·Phase 구조에 슬라이드를 대응시킨다.

배점 대응: 논리적 타당성 30 · 리스크 관리 30 · 기술적 완성도 20 · 성과/발표 20.

실행: python build_deck.py
"""
from __future__ import annotations

import copy
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── 디자인 토큰 ──────────────────────────────────────────────────
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GOLD = RGBColor(0xC8, 0xA2, 0x4B)
INK = RGBColor(0x1A, 0x1F, 0x27)
BODY = RGBColor(0x44, 0x50, 0x5E)
MUTED = RGBColor(0x76, 0x82, 0x90)
LINE = RGBColor(0xD8, 0xDE, 0xE6)
BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB3, 0x3A, 0x3A)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)

FONT = "Malgun Gothic"
W, H = Inches(13.333), Inches(7.5)
ML = Inches(0.85)                      # 좌측 기준선
CW = W - ML * 2                        # 본문 폭

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "wt-freeze" / "_deck"
OUT = Path(r"C:/EST/indi/aq2_qrisproject/doc/헬름_IRP세븐서티액티브_발표자료.pptx")

M = json.loads((ASSETS / "metrics.json").read_text(encoding="utf-8"))


def pct(v: float, dp: int = 1) -> str:
    """음수에 ASCII 하이픈 대신 조판용 마이너스(−)를 쓴다 — 문서 전체 표기 통일."""
    return f"{v:.{dp}f}%".replace("-", "−")


def num(v: float, dp: int = 2) -> str:
    return f"{v:.{dp}f}".replace("-", "−")


V2, V1, BM = M["v2"]["strategy"], M["v1"]["strategy"], M["v2"]["benchmark"]
CUT2, CUT1 = M["cut2025_v2"], M["cut2025_v1"]
C3 = M["cost_3x"]["strategy"]

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ── 저수준 헬퍼 ─────────────────────────────────────────────────
def _ea(run) -> None:
    """한글이 라틴 폰트로 떨어지지 않도록 East-Asian typeface 를 직접 지정한다."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def text(slide, x, y, w, h, s, size=16, color=BODY, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP):
    """텍스트 상자 하나. s 는 문자열 또는 (문자열, 옵션dict) 리스트."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = s if isinstance(s, list) else [s]
    for i, item in enumerate(items):
        txt, opt = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opt.get("align", align)
        p.line_spacing = opt.get("spacing", spacing)
        if opt.get("space_before"):
            p.space_before = Pt(opt["space_before"])
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(opt.get("size", size))
        r.font.bold = opt.get("bold", bold)
        r.font.color.rgb = opt.get("color", color)
        r.font.name = FONT
        _ea(r)
    return tb


def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def slide_base(title=None, kicker=None, rule=True):
    """공통 슬라이드 골격 — 상단 커커·제목·구분선."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=WHITE)
    y = Inches(0.52)
    if kicker:
        text(s, ML, y, CW, Inches(0.28), kicker, size=12, color=GOLD, bold=True)
        y += Inches(0.34)
    if title:
        text(s, ML, y, CW, Inches(0.6), title, size=27, color=NAVY, bold=True)
        y += Inches(0.72)
    if rule:
        rect(s, ML, y, CW, Emu(9525), fill=LINE)
        y += Inches(0.26)
    return s, y


def bullets(slide, x, y, w, items, size=15, gap=0.42, dot=GOLD):
    """· 불릿 목록. items 원소는 문자열 또는 (본문, 강조머리) 튜플."""
    for i, it in enumerate(items):
        yy = y + Inches(gap * i)
        rect(slide, x, yy + Inches(0.09), Inches(0.075), Inches(0.075), fill=dot)
        body = it if isinstance(it, str) else it[0]
        text(slide, x + Inches(0.22), yy, w - Inches(0.22), Inches(gap), body, size=size)
    return y + Inches(gap * len(items))


def table(slide, x, y, w, rows, col_w=None, header=True, size=13, row_h=0.36,
          align=None, emphasis=None):
    """간결한 표. rows[0] 을 헤더로 쓴다. emphasis: 강조할 (행,열) 집합."""
    nr, nc = len(rows), len(rows[0])
    shape = slide.shapes.add_table(nr, nc, x, y, w, Inches(row_h * nr))
    tbl = shape.table
    tbl.first_row = header
    tbl.horz_banding = False
    if col_w:
        total = sum(col_w)
        for j, cwf in enumerate(col_w):
            tbl.columns[j].width = Emu(int(w * cwf / total))
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(row_h)
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = ""
            c.margin_left = c.margin_right = Inches(0.09)
            c.margin_top = c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = NAVY if (header and i == 0) else (
                BG if i % 2 == 0 else WHITE)
            p = c.text_frame.paragraphs[0]
            p.alignment = (PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER) \
                if align is None else align[j]
            r = p.add_run()
            r.text = str(val)
            strong = (header and i == 0) or (emphasis and (i, j) in emphasis)
            r.font.size = Pt(size)
            r.font.bold = strong
            r.font.color.rgb = WHITE if (header and i == 0) else (
                NAVY if strong else BODY)
            r.font.name = FONT
            _ea(r)
    return shape


def kpi_row(slide, y, tiles, h=1.18):
    """숫자 타일 가로 배치 — (값, 라벨, 보조) 튜플."""
    n = len(tiles)
    gap = Inches(0.22)
    tw = Emu(int((CW - gap * (n - 1)) / n))
    for i, (val, label, sub) in enumerate(tiles):
        x = ML + Emu(int((tw + gap) * i))
        rect(slide, x, y, tw, Inches(h), fill=BG)
        rect(slide, x, y, Inches(0.05), Inches(h), fill=GOLD)
        # 라벨·보조문구는 타일 하단에서 역산 배치 — h 를 줄여도 타일 밖으로 넘치지 않게.
        text(slide, x + Inches(0.24), y + Inches(0.15), tw - Inches(0.36), Inches(0.5),
             val, size=min(31, 31 * h / 1.18), color=NAVY, bold=True)
        text(slide, x + Inches(0.24), y + Inches(h - 0.50), tw - Inches(0.36), Inches(0.22),
             label, size=12, color=INK, bold=True)
        if sub:
            text(slide, x + Inches(0.24), y + Inches(h - 0.26), tw - Inches(0.36), Inches(0.2),
                 sub, size=10.5, color=MUTED)


def picture(slide, path, y, height=None, width=None):
    """가운데 정렬 이미지."""
    p = ASSETS / path
    pic = slide.shapes.add_picture(str(p), Inches(0), y,
                                   height=height, width=width)
    pic.left = Emu(int((W - pic.width) / 2))
    return pic


def footer(slide, note):
    text(slide, ML, H - Inches(0.52), CW, Inches(0.3), note, size=10, color=MUTED)


# ══════════════════════════════════════════════════════════════════
# 1. 표지
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, H - Inches(0.13), W, Inches(0.13), fill=GOLD)
text(s, ML, Inches(1.55), CW, Inches(0.4), "AI퀀트 2기 2조 · 팀프로젝트 — 액티브 ETF 개발",
     size=14, color=GOLD, bold=True)
text(s, ML, Inches(2.25), CW, Inches(1.1), "헬름 IRP세븐서티액티브",
     size=52, color=WHITE, bold=True)
text(s, ML, Inches(3.5), CW, Inches(0.9),
     "확신의 크기만큼만 싣는다 —\n채권 30%로 지키고, 70%를 추세 점수에 비례해 싣는 하락 방어형 글로벌 액티브 ETF",
     size=17, color=RGBColor(0xD3, 0xDC, 0xE8), spacing=1.4)
rect(s, ML, Inches(4.85), Inches(2.6), Emu(9525), fill=GOLD)
text(s, ML, Inches(5.15), CW, Inches(1.0),
     [("백테스트 2020.01 ~ 2026.06  ·  벤치마크 KODEX TRF7030", {"size": 13}),
      ("CAGR 13.6%   ·   MDD −12.7%   ·   Calmar 1.07 (벤치 0.60)",
       {"size": 15, "bold": True, "color": WHITE, "space_before": 8})],
     size=13, color=RGBColor(0xA9, 0xB8, 0xCC))
text(s, ML, H - Inches(1.0), CW, Inches(0.3),
     "2026.07  ·  전략·검증 / 엔진·데이터 / 상품·발표", size=12,
     color=RGBColor(0x8C, 0x9D, 0xB5))

# ══════════════════════════════════════════════════════════════════
# 2. 한 장 요약
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("한 장 요약", "EXECUTIVE SUMMARY")
kpi_row(s, y, [
    (pct(V2["cagr_pct"]), "CAGR", f"벤치마크 {pct(BM['cagr_pct'])}"),
    (pct(V2["mdd_pct"]), "최대낙폭 (MDD)", f"벤치마크 {pct(BM['mdd_pct'])}"),
    (num(V2["calmar"]), "Calmar", f"벤치마크 {num(BM['calmar'])}"),
    ("+2.2%", "최악의 해 (2022)", "벤치마크 −11.3%"),
])
text(s, ML, y + Inches(1.30), CW, Inches(0.26),
     "※ 벤치마크 = KODEX TRF7030 — 선진국주식 70 / 국내채권 30 의 동일 체급 실존 IRP 상품 (선정 근거: Phase 1)",
     size=11.5, color=MUTED)
y2 = y + Inches(1.68)
text(s, ML, y2, CW, Inches(0.4),
     "수익은 벤치마크와 비슷하다. 다른 것은 그 수익을 얻기까지 감수한 낙폭이다.",
     size=19, color=NAVY, bold=True)
bullets(s, ML, y2 + Inches(0.62), CW, [
    "같은 목적지에 절반의 멀미로 — CAGR 13.6% vs 13.2%, MDD −12.7% vs −22.1%",
    "위험조정 성과는 1.8배 — Calmar 1.07 vs 0.60. 2026년을 빼도 1.09로 유지된다",
    "7년 중 최악의 해가 +2.2% — 손실 난 해가 한 번도 없다",
    "24개월 이상 보유 시 손실 확률 0% (어느 달에 시작했든, 적립·거치 모두)",
    "방어는 재량이 아니라 규칙이다 — TrendScore 45 하향 시 전량 청산, 전 기간 자동 집행",
], size=15.5, gap=0.44)
footer(s, "거래비용 왕복 0.10% 반영 · 세전 · 룩어헤드 방지 검증 완료 · 재현 브랜치 v2-tier2a-freeze")

# ══════════════════════════════════════════════════════════════════
# 3. Phase 1 — 왜 이 상품인가
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("왜 이 상품인가", "PHASE 1 · 상품기획")
text(s, ML, y, CW, Inches(0.4),
     "퇴직연금 적립금은 커지는데, 가입자 수익률은 정체돼 있다.", size=19,
     color=NAVY, bold=True)
y += Inches(0.66)
cw = Emu(int((CW - Inches(0.3)) / 2))
rect(s, ML, y, cw, Inches(2.5), fill=BG)
text(s, ML + Inches(0.28), y + Inches(0.24), cw - Inches(0.56), Inches(0.3),
     "시장 배경", size=14, color=GOLD, bold=True)
bullets(s, ML + Inches(0.28), y + Inches(0.66), cw - Inches(0.56), [
    "액티브 ETF 급성장 — 비교지수 상관계수 요건 폐지 추진으로 운용 자유도 확대",
    "IRP 가입자는 원리금보장·지수추종에 편중 — 방어 설계된 대안이 드물다",
    "연 900만 세액공제 한도, 55세까지 인출 제한 = 구조적 장기 계좌",
], size=13, gap=0.55)
rect(s, ML + cw + Inches(0.3), y, cw, Inches(2.5), fill=NAVY)
text(s, ML + cw + Inches(0.58), y + Inches(0.24), cw - Inches(0.56), Inches(0.3),
     "우리의 답", size=14, color=GOLD, bold=True)
bullets(s, ML + cw + Inches(0.58), y + Inches(0.66), cw - Inches(0.56), [
    "IRP 계좌를 정면 타깃으로 삼은 70:30 액티브 ETF",
    "상승장 참여 + 하락장 방어를 '재량'이 아니라 '규칙'으로 구현",
    "장기 계좌에는 높은 수익률보다 견딜 수 있는 낙폭이 중요하다",
], size=13, gap=0.55, dot=GOLD)
for sh in s.shapes:
    if sh.has_text_frame and sh.left and sh.left > ML + cw:
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                if r.font.color and r.font.color.rgb == BODY:
                    r.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
y += Inches(2.78)
text(s, ML, y, CW, Inches(0.4),
     "타깃 고객 — 30~40대 IRP 가입자. 세액공제를 챙기되, 2022년 같은 해에 계좌를 열어보기 "
     "두렵지 않기를 바라는 사람.", size=14, color=BODY)

# ══════════════════════════════════════════════════════════════════
# 4. 비교지수 선정 — 이 뒤의 모든 비교가 이 슬라이드에 기댄다
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("비교지수 — 왜 KODEX TRF7030 인가", "PHASE 1 · ① 비교지수 선정")

rect(s, ML, y, CW, Inches(0.94), fill=NAVY)
text(s, ML + Inches(0.3), y + Inches(0.17), CW - Inches(0.6), Inches(0.3),
     "KODEX TRF7030", size=19, color=WHITE, bold=True)
text(s, ML + Inches(0.3), y + Inches(0.53), CW - Inches(0.6), Inches(0.3),
     "TRF = Target Risk Fund · 선진국주식 70% / 국내채권 30% 고정 배분 · IRP 계좌에서 실제 판매 중",
     size=12.5, color=RGBColor(0xC3, 0xCF, 0xDF))
y += Inches(1.22)

reasons = [
    ("체급이 같다",
     "우리도 위험자산 70 / 채권 30 이다. 배분 비율이 같아야 '자산배분의 차이'가 아니라 "
     "'운용의 차이'만 비교된다 — 주식 100% 지수와 비교하면 방어력은 배분 덕인지 운용 덕인지 알 수 없다."),
    ("가상 지수가 아니다",
     "가입자가 지금 IRP 계좌에서 실제로 살 수 있는 대안이다. "
     "\"이 상품 대신 무엇을 사겠는가\"에 대한 정직한 답을 비교 대상으로 삼았다."),
    ("규제상 자유 비교가 가능하다",
     "완전 자율 액티브로 설계해 비교지수 상관계수 요건을 적용받지 않는다. "
     "따라서 지수를 따라갈 의무 없이, 이길 대상으로만 삼는다."),
]
for i, (head, desc) in enumerate(reasons):
    yy = y + Inches(0.92 * i)
    text(s, ML, yy + Inches(0.02), Inches(0.5), Inches(0.4), f"{i+1}",
         size=22, color=GOLD, bold=True)
    text(s, ML + Inches(0.5), yy, CW - Inches(0.5), Inches(0.28), head,
         size=16, color=NAVY, bold=True)
    text(s, ML + Inches(0.5), yy + Inches(0.33), CW - Inches(0.5), Inches(0.5),
         desc, size=12.5, color=BODY, spacing=1.3)

y += Inches(2.88)
rect(s, ML, y, CW, Inches(1.0), fill=BG)
rect(s, ML, y, Inches(0.05), Inches(1.0), fill=GOLD)
text(s, ML + Inches(0.26), y + Inches(0.15), CW - Inches(0.55), Inches(0.28),
     "그래서 무엇이 달라지나 — 같은 70:30, 다른 운용", size=14, color=NAVY, bold=True)
text(s, ML + Inches(0.26), y + Inches(0.48), CW - Inches(0.55), Inches(0.44),
     "벤치마크는 70%를 언제나 만충으로 들고 간다. 본 상품은 같은 70%를 추세 점수에 비례해 채운다"
     "(평균 49.6%).\n수익률이 비슷하게 나오는 것은 당연하고, 차이는 낙폭에서 벌어진다 — 이 비교의 핵심은 그 지점이다.",
     size=12.5, color=BODY, spacing=1.3)

# ══════════════════════════════════════════════════════════════════
# 5. 투자 철학
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("투자 철학 — 네 문장", "PHASE 1 · 상품컨셉")
philos = [
    ("규율이 곧 방어다", "채권 30% 상시 고정. 어떤 국면에도 하방 한계를 구조화한다."),
    ("확신의 크기만큼만 싣는다",
     "자격선(52점)을 넘은 종목만, 점수에 비례해 슬롯의 30~100%를 충전한다.\n"
     "살지 말지(스위치)가 아니라 얼마나 살지(조광기)를 정한다."),
    ("살 만한 것이 없으면 사지 않는다",
     "기준 미달이면 그 자리는 채권으로. 77회 체크 중 만충은 18회뿐, 평균 노출은 만충의 71%."),
    ("무리하지 않는다", "무레버리지·무파생. 연금계좌 매수 가능 요건을 스스로 충족한다."),
]
for i, (head, desc) in enumerate(philos):
    yy = y + Inches(1.28 * i)
    text(s, ML, yy + Inches(0.02), Inches(0.55), Inches(0.5), f"{i+1}",
         size=30, color=GOLD, bold=True)
    text(s, ML + Inches(0.62), yy, CW - Inches(0.62), Inches(0.32), head,
         size=18, color=NAVY, bold=True)
    text(s, ML + Inches(0.62), yy + Inches(0.42), CW - Inches(0.62), Inches(0.7),
         desc, size=14, color=BODY, spacing=1.3)

# ══════════════════════════════════════════════════════════════════
# 6. TrendScore — 국면 판단 (논리 배점의 핵심)
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("TrendScore — 국면을 읽는 커스텀 지표", "PHASE 2-A · 시장 국면 판단")
text(s, ML, y, CW, Inches(0.4),
     "\"이 종목이 지금 믿을 만한 추세 위에 있는가\"를 0~100의 한 숫자로 답한다.",
     size=18, color=NAVY, bold=True)
y += Inches(0.52)
text(s, ML, y, CW, Inches(0.62),
     "지표 하나에는 실패 국면이 하나씩 따라붙는다 — 이동평균은 횡보장에서 속임수 신호를 내고, "
     "모멘텀은 고점에서 가장 강하게 사라고 한다.\nTrendScore는 서로의 맹점을 가리는 네 팩터를 겹쳐, "
     "한 지표가 틀리는 국면에서 나머지가 점수를 끌어내린다.",
     size=13.5, color=BODY, spacing=1.3)
y += Inches(0.86)
table(s, ML, y, CW, [
    ["팩터", "비중", "무엇을 보는가", "이 팩터가 막는 실패"],
    ["EWMAC 앙상블", "55%", "추세의 방향과 지속성 (8/32·16/64·32/128일 3쌍)", "특정 기간 설정 과최적화"],
    ["TSMOM (12−1)", "25%", "12개월 수익률 − 직전 1개월", "단기 급등 추격매수"],
    ["RSI (14일)", "20%", "최근 2~3주 매수·매도 압력 균형", "중·장기 지표의 반응 지연"],
    ["ADX 페널티", "−15점", "추세의 강도 (방향이 아님)", "횡보장 오탐"],
], col_w=[1.0, 0.45, 1.85, 1.3], size=12.5, row_h=0.42,
    align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
y += Inches(2.32)
rect(s, ML, y, CW, Inches(1.12), fill=BG)
rect(s, ML, y, Inches(0.05), Inches(1.12), fill=GOLD)
text(s, ML + Inches(0.26), y + Inches(0.16), CW - Inches(0.55), Inches(0.28),
     "변동성 정규화 — 이 지표의 핵심 장치", size=14, color=NAVY, bold=True)
text(s, ML + Inches(0.26), y + Inches(0.5), CW - Inches(0.55), Inches(0.55),
     "모든 팩터를 연율 변동성으로 나눈다. 재는 것은 '추세의 크기'가 아니라 '추세의 명료성' —\n"
     "그래서 반도체와 금현물처럼 성격이 다른 36종을 같은 자로 재서 한 줄로 세울 수 있다.",
     size=12.5, color=BODY, spacing=1.3)

# ══════════════════════════════════════════════════════════════════
# 7. 유니버스
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("투자 유니버스 — 36종", "PHASE 2-B · 종목선별")
table(s, ML, y, CW, [
    ["슬리브", "구성", "역할"],
    ["한국 섹터 ETF 15종", "IT·반도체·증권·건설 등", "국내 주도 테마 포착"],
    ["해외지수·미국섹터", "S&P500·나스닥100·S&P 섹터·일본·중국·인도", "글로벌 성장 참여"],
    ["원자재·리츠", "금현물(411060)·구리실물·리츠인프라", "체제 보험·분산"],
    ["채권 3종 (고정 30%)", "단기채·국고채3년·종합채권액티브", "구조적 방어벽"],
], col_w=[1.15, 2.1, 1.35], size=13.5, row_h=0.46,
    align=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
y += Inches(2.55)
text(s, ML, y, CW, Inches(0.35),
     "유니버스는 종목 명단이 아니라 자격 규정이다", size=17, color=NAVY, bold=True)
bullets(s, ML, y + Inches(0.5), CW, [
    "국내 상장 · 무레버리지 · 연금계좌 편입 가능 · 상장 후 252거래일 경과(지표 워밍업) · 자산군 대표성",
    "신규 상장 ETF는 252거래일 축적 후 재심사 편입 — 백테스트와 실전은 같은 규정의 다른 시점 단면이다",
    "원자재는 백테스트 6.5년에 없던 국면(인플레 재점화 등)에 대한 체제 보험 — 그 비용을 측정·공시한다",
], size=14, gap=0.42)

# ══════════════════════════════════════════════════════════════════
# 7. 운용 프로세스
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("운용 프로세스 — 3단계", "PHASE 2-B · 비중 결정 · 리밸런싱")
steps = [
    ("STEP 1", "종목별 추세 점수화 (월간)",
     "36종 각각에 TrendScore(0~100) 산출.\n변동성 정규화로 자산군이 달라도 한 줄로 비교한다."),
    ("STEP 2", "선별과 충전 (핵심 규칙)",
     "점수 상위 최대 7종목. 52점 이상만 자격,\n점수에 비례해 충전 — 52점=슬롯 30% → 60점 이상=만충.\n"
     "45점 미만 시 전량 매도. 빈 슬롯은 단기채 대피."),
    ("STEP 3", "리밸런싱",
     "분기 정기 + 목표 대비 ±7%p 이탈 시 수시.\n월 적립 현금흐름이 자연 리밸런싱을 보조."),
]
cw3 = Emu(int((CW - Inches(0.44)) / 3))
for i, (tag, head, desc) in enumerate(steps):
    x = ML + Emu(int((cw3 + Inches(0.22)) * i))
    rect(s, x, y, cw3, Inches(2.62), fill=BG)
    rect(s, x, y, cw3, Inches(0.06), fill=GOLD)
    text(s, x + Inches(0.26), y + Inches(0.28), cw3 - Inches(0.52), Inches(0.26),
         tag, size=12, color=GOLD, bold=True)
    text(s, x + Inches(0.26), y + Inches(0.60), cw3 - Inches(0.52), Inches(0.5),
         head, size=16, color=NAVY, bold=True)
    text(s, x + Inches(0.26), y + Inches(1.22), cw3 - Inches(0.52), Inches(1.2),
         desc, size=12.5, color=BODY, spacing=1.35)
y += Inches(2.95)
text(s, ML, y, CW, Inches(0.35),
     "진입 52 / 청산 45 — 두 문턱의 간격이 완충지대를 만든다", size=16, color=NAVY, bold=True)
text(s, ML, y + Inches(0.42), CW, Inches(0.6),
     "52점 이상을 유지하는 보유는 웬만한 신규 후보에 흔들리지 않는다. 문턱을 하나만 두면 그 근처에서 "
     "사고팔기를 반복(휩쏘)하게 되므로, 들어오는 문과 나가는 문을 다르게 뒀다.",
     size=13.5, color=BODY, spacing=1.3)

# ══════════════════════════════════════════════════════════════════
# 8. 점수 비례 충전 + 노출 차트
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("노출은 스스로 내려간다", "운용 결과 · 실효 노출률 실측")
text(s, ML, y, CW, Inches(0.34),
     "하락장이 추세를 꺾으면 점수가 내려가고, 점수가 내려가면 노출이 줄어든다.",
     size=17, color=NAVY, bold=True)
text(s, ML, y + Inches(0.36), CW, Inches(0.28),
     "판단이 아니라 인과다 — 그래서 다음 위기에도 같은 방식으로 작동한다.",
     size=13, color=BODY)
picture(s, "exposure.png", y + Inches(0.76), height=Inches(2.78))
y2 = y + Inches(3.66)
kpi_row(s, y2, [
    (pct(M["exposure"]["mean_pct"]), "평균 실효 노출", "만충(70%)의 71%"),
    ("0%", "최저 노출 (2020.05)", "코로나 국면 전량 대피"),
    ("18 / 77", "만충 개월 수", "나머지는 부분 충전 또는 대피"),
    ("6.1종", "평균 보유 종목", "종목당 평균 8.1%"),
], h=1.18)

# ══════════════════════════════════════════════════════════════════
# 9. 리스크 관리 (배점 30)
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("리스크 관리 체계 — 룰과 실측", "RISK MANAGEMENT · 명시적 룰")
table(s, ML, y, CW, [
    ["요구 룰", "본 상품의 구현", "백테스트 실측"],
    ["손절 기준", "TrendScore 45 하향 시 전량 청산 — 가격이 아닌 추세 소멸 기준",
     "2020.5·2023.5 전량 대피 포함 전 기간 자동 집행"],
    ["비중 조절 룰", "점수 비례 충전: 52점=30% → 60점=만충, 식으면 자동 감량",
     "평균 노출 49.6%, 위기 시 한 자릿수"],
    ["MDD 통제", "−15% 서킷브레이커: 도달 시 위험자산 절반 축소 후 월간 재심사",
     "전 기간(MDD −12.7%) 무발동 — 최후 방어선"],
    ["개별 종목 한도", "위험자산 종목당 최대 10%", "동일가중 상한으로 강제, 위반 0회"],
    ["현금성 상한", "현금성(단기채 대피분)은 사테라이트 70% 이내로 구조적 제한",
     "대피 최대치 = 슬롯 전체 공실 시 70%"],
    ["집중 완화", "자산군 분산 + Top-7 분산", "평균 보유 6.1종목, 종목당 평균 8.1%"],
], col_w=[0.75, 2.1, 1.55], size=11.5, row_h=0.5,
    align=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
y += Inches(3.75)
rect(s, ML, y, CW, Inches(0.86), fill=BG)
rect(s, ML, y, Inches(0.05), Inches(0.86), fill=GOLD)
text(s, ML + Inches(0.26), y + Inches(0.13), CW - Inches(0.5), Inches(0.28),
     "서킷브레이커는 한 번도 발동하지 않았다", size=14, color=NAVY, bold=True)
text(s, ML + Inches(0.26), y + Inches(0.44), CW - Inches(0.5), Inches(0.34),
     "게이트와 경사가 그보다 앞서 노출을 줄이기 때문이다. 룰은 그 1차 방어가 실패할 때를 위한 최후 방어선으로 존재한다.",
     size=12.5, color=BODY)

# ══════════════════════════════════════════════════════════════════
# 10. 낙폭 — 방어 실증
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("방어의 실증 — 같은 시장에서 얼마나 덜 잃었나", "RISK MANAGEMENT · 방어력 실측")
picture(s, "drawdown.png", y + Inches(0.05), height=Inches(3.3))
y2 = y + Inches(3.55)
kpi_row(s, y2, [
    (pct(V2["mdd_pct"]), "본 상품 MDD", f"벤치마크 {pct(BM['mdd_pct'])}"),
    ("+2.2%", "2022년 수익", "벤치마크 −11.3%"),
    ("64%", "구간 승률", "손익비 1.82 · PF 3.23"),
    ("0회", "서킷브레이커 발동", "−15% 미도달 — 최후 방어선"),
], h=1.12)
footer(s, "낙폭(drawdown) = 직전 고점 대비 하락률. 값이 0에 가까울수록 방어적이다.")

# ══════════════════════════════════════════════════════════════════
# 11. 성과 — 자산곡선
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("성과 — 자산곡선", "PHASE 3 · 백테스트")
picture(s, "equity.png", y, height=Inches(3.0))
y2 = y + Inches(3.2)
table(s, ML, y2, CW, [
    ["", "본 상품", "KODEX TRF7030", "차이"],
    ["CAGR", pct(V2["cagr_pct"]), pct(BM["cagr_pct"]), "＋0.4%p"],
    ["MDD", pct(V2["mdd_pct"]), pct(BM["mdd_pct"]), "낙폭 절반"],
    ["Sharpe", num(V2["sharpe"]), num(BM["sharpe"]), "＋0.02"],
    ["Calmar", num(V2["calmar"]), num(BM["calmar"]), "1.8배"],
    ["2026 제외 Calmar", num(CUT2["strategy"]["calmar"]),
     num(CUT2["benchmark"]["calmar"]), "1.9배"],
], col_w=[1.3, 1.0, 1.2, 0.9], size=12, row_h=0.31,
    emphasis={(i, 1) for i in range(1, 6)})
footer(s, "거래비용 왕복 0.10% 반영 · 세전 · 비용을 3배(0.30%)로 올려도 MDD −12.7% 불변, Calmar 1.01")

# ══════════════════════════════════════════════════════════════════
# 12. 연도별
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("연도별 성과 — 2022년이 이 상품의 정체성", "PHASE 3 · 국면별 방어력 분해")
picture(s, "yearly.png", y, height=Inches(2.8))
y2 = y + Inches(3.02)


def sgn(v: float) -> str:
    """부호 포함 표기 + 조판용 마이너스."""
    return f"{v:+.1f}".replace("-", "−")


ys_ = M["v2"]["yearly"]
rows = [["연도"] + [str(r["year"]) if r["year"] != 2026 else "2026상" for r in ys_]]
rows.append(["본 상품"] + [sgn(r["strat"]) for r in ys_])
rows.append(["벤치마크"] + [sgn(r["bench"]) for r in ys_])
rows.append(["초과"] + [sgn(r["excess"]) for r in ys_])
table(s, ML, y2, CW, rows, size=12, row_h=0.34,
      col_w=[1.1] + [0.72] * len(ys_),
      emphasis={(1, 3), (2, 3), (3, 3), (1, 7), (3, 7)})
y2 += Inches(1.52)
text(s, ML, y2, CW, Inches(0.66),
     "7년 중 4년은 벤치마크보다 덜 벌었다. 그러나 벤치마크가 −11.3%였던 2022년에 +2.2%를 냈다.\n"
     "평시에 보험료(소폭 열위)를 내고 위기에 보험금을 받는 구조 — 이기는 해는 적지만, 잃는 해가 없다.",
     size=14, color=BODY, spacing=1.35)

# ══════════════════════════════════════════════════════════════════
# 13. 구현 — 코드 구현력(기술 배점의 절반)
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("구현 — 계층 분리와 세 개의 안전장치", "PHASE 3 · 기술 구현")
text(s, ML, y, CW, Inches(0.36),
     "재현되지 않는 백테스트는 결과가 아니라 주장이다.",
     size=18, color=NAVY, bold=True)
text(s, ML, y + Inches(0.40), CW, Inches(0.3),
     "숫자를 만드는 코드보다, 그 숫자가 틀렸을 때 시끄럽게 실패하는 코드에 더 공을 들였다.",
     size=13, color=BODY)
y += Inches(0.88)

guards = [
    ("① 룩어헤드 원천 차단",
     "신호는 봉 i 종가로 확정하고 **체결은 i+1 시가**로 한다. 전략이 오늘 종가를 보고 "
     "오늘 사는 일이 코드 구조상 불가능하다 — 규칙이 아니라 엔진이 막는다."),
    ("② 무결성 fail-loud",
     "실제 사고에서 나온 장치다. 금(411060) 데이터가 없는 환경에서 돌리면 경고 한 줄만 남고 "
     "금이 빠진 채 그럴듯한 수치가 나왔다. 이제 실행 전 전 종목 로드를 검증하고 "
     "하나라도 실패하면 **멈춘다.**"),
    ("③ 관측 프로브 무해성",
     "노출·보유 구성을 재려고 판정 로직을 다시 구현하면 엔진과 어긋난다. 프로브는 엔진이 "
     "**이미 정한** 값을 받아적기만 하며, 프로브를 끼운 자산곡선이 끼우지 않은 것과 "
     "**완전히 일치함**을 매 실행 검증한다."),
]
cw3 = Emu(int((CW - Inches(0.44)) / 3))
for i, (head, desc) in enumerate(guards):
    x = ML + Emu(int((cw3 + Inches(0.22)) * i))
    rect(s, x, y, cw3, Inches(2.30), fill=BG)
    rect(s, x, y, cw3, Inches(0.05), fill=GOLD)
    text(s, x + Inches(0.24), y + Inches(0.26), cw3 - Inches(0.48), Inches(0.5),
         head, size=14, color=NAVY, bold=True)
    # 카드 본문은 **강조** 마크업을 쓰지 않으므로 별표를 제거해 넘긴다
    text(s, x + Inches(0.24), y + Inches(0.74), cw3 - Inches(0.48), Inches(1.4),
         desc.replace("**", ""), size=11.5, color=BODY, spacing=1.32)

y += Inches(2.62)
rect(s, ML, y, CW, Inches(1.06), fill=NAVY)
text(s, ML + Inches(0.28), y + Inches(0.16), CW - Inches(0.56), Inches(0.28),
     "아키텍처 · 스택", size=12, color=GOLD, bold=True)
text(s, ML + Inches(0.28), y + Inches(0.46), CW - Inches(0.56), Inches(0.5),
     "책임 하나 = 클래스 하나 = 파일 하나 · 계층 9개(data · indicator · strategy · backtest · "
     "satellite · irp · portfolio · analysis · report) · 파이썬 67파일 8,123줄\n"
     "추상 기반 클래스로 인터페이스 고정 + 의존성 주입 → 지표·전략 교체가 설정 변경으로 끝난다  ·  "
     "pandas · numpy · pyarrow(Parquet) · matplotlib",
     size=11, color=RGBColor(0xD3, 0xDC, 0xE8), spacing=1.35)

# ══════════════════════════════════════════════════════════════════
# 14. 검증 (기술 배점 20)
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("왜 이 숫자를 믿을 수 있나", "PHASE 3 · 검증과 재현성")
items = [
    ("채점 기준 사전 고정", "두 구간 Calmar·최악 해·2022 방어를 결과 확인 전에 확정. "
     "미달 실험(진입 경사 단독 MDD −17.6% 등)은 전부 기각·기록."),
    ("9/9 plateau", "채택값 주변 9개 설정을 전수 측정, 전부 기준선 초과 — 한 점의 운이 아니라 면 전체가 들렸다."),
    ("두 구간 교차 검증", f"2026년이라는 한 해를 통째로 빼도 Calmar {num(CUT2['strategy']['calmar'])}로 유지·상승한다"
     f"(전 구간 {num(V2['calmar'])}). 같은 구간 벤치마크는 {num(CUT2['benchmark']['calmar'])} — 특정 국면에 기댄 성과가 아니다."),
    ("개선 전수 탐색 후 확정", "동결 이후 대안 8계열·11개 구성을 같은 기준으로 전수 측정 — 채택 0. "
     "\"대안을 못 찾아서\"가 아니라 \"주변 전부를 재본 뒤\" 확정된 구성이다."),
    ("재현성", "동일 커밋·동일 데이터에서 자산곡선이 그대로 재현된다. 기능 플래그 대조 실험으로 "
     "각 규칙의 기여를 개별 검증했고, 검증 32항목과 자가 발견 결함 전수를 문서로 공개한다."),
]
for i, (head, desc) in enumerate(items):
    yy = y + Inches(0.98 * i)
    rect(s, ML, yy + Inches(0.06), Inches(0.05), Inches(0.62), fill=GOLD)
    text(s, ML + Inches(0.24), yy, CW - Inches(0.24), Inches(0.3), head,
         size=15.5, color=NAVY, bold=True)
    text(s, ML + Inches(0.24), yy + Inches(0.33), CW - Inches(0.24), Inches(0.55),
         desc, size=13, color=BODY, spacing=1.3)

# ══════════════════════════════════════════════════════════════════
# 14. 가입자 시나리오
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("가입자 시나리오 — 연 900만 (세액공제 한도)", "가입자 관점 · 적립 시뮬레이션")
table(s, ML, y, CW, [
    ["납입 방식", "총 납입", "최종 평가액", "손익률", "MWR 연율"],
    ["월 75만 적립식", "5,850만", "1억 394만  (벤치 9,790만)", "+77.7%", "17.5%  (벤치 15.7%)"],
    ["연초 900만 거치식", "6,300만", "1억 1,469만  (벤치 1억 725만)", "+82.0%", "17.1%  (벤치 15.2%)"],
], col_w=[1.15, 0.8, 1.75, 0.75, 1.15], size=13, row_h=0.42,
    align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
y += Inches(1.45)
kpi_row(s, y, [
    ("0%", "24개월 보유 시 손실 확률", "적립·거치 모두 · 시작 월 무관"),
    ("7.6%", "12개월 적립 손실 확률", "벤치마크 13.6%"),
    ("3.7%", "벤치마크 24개월 손실 확률", "본 상품은 0%"),
], h=1.05)
y += Inches(1.28)
text(s, ML, y, CW, Inches(0.32),
     "12개월 단기는 손실이 가능하다 — 단기 자금용이 아니며, 55세까지 묶이는 IRP의 성격과 정합한다.",
     size=13.5, color=BODY)
y += Inches(0.44)
text(s, ML, y, CW, Inches(0.3), "10년 예상 적립금 (가정 기반 산수 예시 — 백테스트 실적 아님)",
     size=13, color=NAVY, bold=True)
table(s, ML, y + Inches(0.37), Inches(7.4), [
    ["연 수익률 가정", "월 75만 적립 (120회)", "연초 900만 거치 (10회)"],
    ["6% (보수적)", "1.22억", "1.26억"],
    ["10% (중립)", "1.51억", "1.58억"],
    ["13.6% (백테스트 CAGR)", "1.83억", "1.94억"],
], col_w=[1.2, 1.0, 1.0], size=11.5, row_h=0.3)
footer(s, "납입원금 9,000만 기준 · 세전 · 보수·세금 미반영 · 수익률이 매년 일정하다는 단순 복리 가정")

# ══════════════════════════════════════════════════════════════════
# 15. 유의사항
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("유의사항 — 정직한 고지", "RISK DISCLOSURE")
text(s, ML, y, CW, Inches(0.4),
     "이 상품이 못하는 것을 먼저 말한다.", size=18, color=NAVY, bold=True)
y += Inches(0.62)
warns = [
    "원금 비보장 실적배당형이며, 백테스트는 미래 수익을 보장하지 않는다.",
    "7년 중 4년은 벤치마크보다 덜 벌었다. 매년 이기는 상품이 아니라, "
    "평시에 보험료를 내고 위기에 보험금을 받는 상품이다.",
    "급반전(V자) 국면에 구조적으로 후행한다 — 추세 확인 후 진입하므로 반등 초입 일부를 포기한다. "
    "2023년 벤치마크 대비 −14.1%p 가 이 약점의 최대 실측치다.",
    "백테스트는 6.5년(해외 ETF 상장 시점 한계)이며, 유니버스는 현존 종목 기준이라 생존편향 가능성이 있다.",
    "세금·슬리피지 미반영(세전).",
    "규제 정합 설계 — 위험자산 비중을 상시 70% 이하로 스스로 제한한다. IRP 계좌의 위험자산 70% 한도와 "
    "같은 선을 상품 내부에 둬, 가입자가 계좌를 이 상품 하나로 채워도 한도를 넘지 않는다.",
]
for i, wtxt in enumerate(warns):
    yy = y + Inches(0.72 * i)
    text(s, ML, yy, Inches(0.4), Inches(0.3), f"{i+1}", size=15, color=GOLD, bold=True)
    text(s, ML + Inches(0.42), yy, CW - Inches(0.42), Inches(0.62), wtxt,
         size=13.5, color=BODY, spacing=1.3)

# ══════════════════════════════════════════════════════════════════
# 16. 상품 개요
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("상품 개요", "PRODUCT SUMMARY")
table(s, ML, y, CW, [
    ["항목", "내용"],
    ["펀드명", "헬름 IRP세븐서티액티브"],
    ["비교지수", "KODEX TRF7030 (선진국주식 70 / 국내채권 30 — 동일 체급 실존 상품)"],
    ["투자 대상", "국내 상장 ETF 36종 (한국 섹터·해외지수·원자재·채권)"],
    ["운용 방식", "추세 점수 비례 충전 + Top-7 선별 + 채권 30% 고정"],
    ["리밸런싱", "월간 점수 체크 · 분기 정기 + 목표 대비 ±7%p 이탈 시 수시"],
    ["총보수 (가정)", "연 0.49% (국내 자산배분형 액티브 ETF 관행 수준)"],
    ["위험등급", "4등급(보통위험) 가정 — 비교지수(4등급)와 동일 체급·동일 법적 분류"],
], col_w=[0.9, 3.4], size=13, row_h=0.46,
    align=[PP_ALIGN.LEFT, PP_ALIGN.LEFT])

# ══════════════════════════════════════════════════════════════════
# 18. 팀 역할 분담 — ETF 운용본부처럼 세 역할이 한 상품을 만든다
# ══════════════════════════════════════════════════════════════════
s, y = slide_base("팀 역할 분담 — 세 역할이 한 상품을 만든다", "TEAM · R&R")
text(s, ML, y, CW, Inches(0.34),
     "ETF 운용본부의 역할 구조를 그대로 나눴고, 각자의 산출물로 결과를 확인할 수 있게 했다.",
     size=15, color=BODY)
y += Inches(0.6)

roles = [
    ("CIO", "최고투자책임자", "전략 · 검증",
     ["상품 컨셉과 비교지수 선정(KODEX TRF7030)",
      "TrendScore 4팩터 설계 및 가중 결정",
      "자산배분 70:30 · 문턱 52/60/45 확정",
      "리스크 룰(손절·서킷브레이커·집중한도) 수립"],
     "채점 기준 사전 고정 · 9/9 plateau 전수 측정 · 대안 11개 구성 기각 기록"),
    ("CTO", "최고기술책임자", "엔진 · 데이터",
     ["데이터 수집·전처리 파이프라인(Parquet 캐시)",
      "백테스트 엔진 및 지표 코드화(67파일 8,123줄)",
      "무결성 fail-loud 가드 · 관측 프로브 구현",
      "재현 브랜치 동결 및 실험 스크립트 분리"],
     "룩어헤드 구조적 차단 · 결정론적 재현 · 검증 32항목 문서화"),
    ("CMO", "최고마케팅책임자", "상품 · 발표",
     ["타깃 고객 정의 및 네이밍(헬름 IRP세븐서티액티브)",
      "투자제안서 v1.3 작성",
      "투자설명서 19p · 간이투자설명서 5p 제작",
      "발표자료 20장 및 발표 대본 구성"],
     "성과 스토리텔링 · 약점 선공개 원칙 · 예상질문 대응 준비"),
]
cw3 = Emu(int((CW - Inches(0.44)) / 3))
for i, (code, title, tag, items, note) in enumerate(roles):
    x = ML + Emu(int((cw3 + Inches(0.22)) * i))
    rect(s, x, y, cw3, Inches(3.94), fill=BG)
    rect(s, x, y, cw3, Inches(0.72), fill=NAVY)
    text(s, x + Inches(0.24), y + Inches(0.10), cw3 - Inches(0.48), Inches(0.32),
         code, size=19, color=WHITE, bold=True)
    text(s, x + Inches(0.24), y + Inches(0.44), cw3 - Inches(0.48), Inches(0.24),
         f"{title} · {tag}", size=10.5, color=GOLD)
    for j, it in enumerate(items):
        yy = y + Inches(0.92 + 0.58 * j)
        rect(s, x + Inches(0.24), yy + Inches(0.07), Inches(0.06), Inches(0.06), fill=GOLD)
        text(s, x + Inches(0.4), yy, cw3 - Inches(0.66), Inches(0.55), it,
             size=11, color=BODY, spacing=1.3)
    rect(s, x + Inches(0.24), y + Inches(3.24), cw3 - Inches(0.48), Emu(9525), fill=LINE)
    text(s, x + Inches(0.24), y + Inches(3.38), cw3 - Inches(0.48), Inches(0.48),
         note, size=10, color=MUTED, spacing=1.28)

y += Inches(4.18)
text(s, ML, y, CW, Inches(0.32),
     "세 역할이 같은 숫자를 쓴다 — 제안서·투자설명서·발표자료가 모두 재현 브랜치 "
     "v2-tier2a-freeze 의 한 벌(metrics.json)에서 나온다.",
     size=12.5, color=BODY)

# ══════════════════════════════════════════════════════════════════
# 17. 마무리
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, H - Inches(0.13), W, Inches(0.13), fill=GOLD)
text(s, ML, Inches(2.35), CW, Inches(0.9), "확신의 크기만큼만 싣는다",
     size=42, color=WHITE, bold=True)
rect(s, ML, Inches(3.5), Inches(2.6), Emu(9525), fill=GOLD)
text(s, ML, Inches(3.85), CW - Inches(1.5), Inches(1.4),
     "시장을 이기는 상품이 아니라, 시장에 남아 있게 하는 상품입니다.\n"
     "55세까지 묶이는 계좌에 필요한 것은 가장 높은 수익률이 아니라 견딜 수 있는 낙폭입니다.",
     size=17, color=RGBColor(0xD3, 0xDC, 0xE8), spacing=1.45)
text(s, ML, Inches(5.6), CW, Inches(0.4),
     "CAGR 13.6%   ·   MDD −12.7%   ·   Calmar 1.07   ·   최악의 해 +2.2%",
     size=16, color=GOLD, bold=True)
text(s, ML, H - Inches(1.0), CW, Inches(0.3),
     "헬름 IRP세븐서티액티브  ·  AI퀀트 2기 2조  ·  재현 브랜치 v2-tier2a-freeze",
     size=12, color=RGBColor(0x8C, 0x9D, 0xB5))

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"저장 완료: {OUT}  (슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장)")
